# database_operations/file_processor.py

import os
import re
import hashlib
import mimetypes
from datetime import datetime
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path

class FileProcessor:
    """Process various file types and extract metadata and content"""
    
    # Supported file extensions
    SUPPORTED_EXTENSIONS = {
        'text': ['.txt', '.md', '.log', '.csv'],
        'pdf': ['.pdf'],
        'doc': ['.doc', '.docx'],
        'excel': ['.xlsx', '.xls'],
        'powerpoint': ['.pptx', '.ppt'],
        'kml': ['.kml', '.kmz'],
        'image': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff'],
        'other': ['.xml', '.json', '.html', '.htm']
    }
    
    # MGRS coordinate patterns
    MGRS_PATTERNS = [
        r'\b\d{1,2}[A-Z]{3}\d{4,10}\b',  # Standard MGRS
        r'\b\d{1,2}\s[A-Z]{3}\s\d{4,10}\b',  # MGRS with spaces
        r'\b\d{1,2}[A-Z]\s[A-Z]{2}\s\d{4,10}\b'  # MGRS with zone spaces
    ]
    
    # GPS coordinate patterns
    GPS_PATTERNS = [
        r'[-+]?\d{1,3}\.\d+[°]?\s*[NS]?\s*,?\s*[-+]?\d{1,3}\.\d+[°]?\s*[EW]?',  # Decimal degrees
        r'\d{1,3}[°]\s*\d{1,2}[\']\s*\d{1,2}[\"]\s*[NS]\s*,?\s*\d{1,3}[°]\s*\d{1,2}[\']\s*\d{1,2}[\"]\s*[EW]'  # DMS
    ]
    
    def __init__(self):
        self.stats = {
            'files_processed': 0,
            'files_failed': 0,
            'coordinates_found': 0,
            'total_size': 0
        }
    
    def get_supported_formats(self) -> List[str]:
        """Return list of all supported file extensions"""
        formats = []
        for category in self.SUPPORTED_EXTENSIONS.values():
            formats.extend([ext.lstrip('.') for ext in category])
        return sorted(formats)
    
    def is_supported_file(self, file_path: str, allowed_types: List[str] = None) -> bool:
        """Check if file type is supported"""
        ext = Path(file_path).suffix.lower()
        
        if allowed_types:
            return ext.lstrip('.') in allowed_types
        
        for category in self.SUPPORTED_EXTENSIONS.values():
            if ext in category:
                return True
        return False
    
    def extract_coordinates(self, text: str) -> Tuple[List[str], List[str]]:
        """Extract MGRS and GPS coordinates from text"""
        mgrs_coords = []
        gps_coords = []
        
        # Extract MGRS coordinates
        for pattern in self.MGRS_PATTERNS:
            matches = re.findall(pattern, text, re.IGNORECASE)
            mgrs_coords.extend(matches)
        
        # Extract GPS coordinates
        for pattern in self.GPS_PATTERNS:
            matches = re.findall(pattern, text, re.IGNORECASE)
            gps_coords.extend(matches)
        
        # Remove duplicates and clean up
        mgrs_coords = list(set([coord.strip().upper() for coord in mgrs_coords]))
        gps_coords = list(set([coord.strip() for coord in gps_coords]))
        
        return mgrs_coords, gps_coords
    
    def extract_keywords(self, text: str) -> List[str]:
        """Extract potential keywords from text"""
        # Simple keyword extraction - can be enhanced
        words = re.findall(r'\b[A-Za-z]{3,}\b', text)
        
        # Filter out common words and get unique terms
        common_words = {'the', 'and', 'for', 'are', 'but', 'not', 'you', 'all', 'can', 'had', 'her', 'was', 'one', 'our', 'out', 'day', 'get', 'has', 'him', 'his', 'how', 'its', 'may', 'new', 'now', 'old', 'see', 'two', 'who', 'boy', 'did', 'she', 'use', 'her', 'way', 'many', 'then', 'them', 'well', 'were', 'been', 'good', 'much', 'some', 'time', 'very', 'when', 'come', 'here', 'just', 'like', 'long', 'make', 'over', 'such', 'take', 'than', 'only', 'little', 'state', 'years', 'people', 'after', 'first', 'never', 'these', 'think', 'where', 'being', 'every', 'great', 'might', 'shall', 'still', 'those', 'under', 'while', 'could', 'other', 'after', 'first', 'never', 'these', 'think', 'where', 'being', 'every', 'great', 'might', 'shall', 'still', 'those', 'under', 'while'}
        
        keywords = [word.lower() for word in words if len(word) > 3 and word.lower() not in common_words]
        
        # Return most frequent keywords (up to 20)
        from collections import Counter
        counter = Counter(keywords)
        return [word for word, count in counter.most_common(20)]
    
    def process_text_file(self, file_path: str) -> Dict[str, Any]:
        """Process plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            mgrs_coords, gps_coords = self.extract_coordinates(content)
            keywords = self.extract_keywords(content)
            
            return {
                'full_text': content[:10000],  # Limit to first 10k chars
                'mgrs_coordinates': mgrs_coords,
                'gps_coordinates': gps_coords,
                'keywords': ' '.join(keywords),
                'file_size': os.path.getsize(file_path)
            }
        except Exception as e:
            raise Exception(f"Error processing text file: {str(e)}")
    
    def process_kml_file(self, file_path: str) -> Dict[str, Any]:
        """Process KML/KMZ files"""
        try:
            import xml.etree.ElementTree as ET
            
            content = ""
            coordinates = []
            
            if file_path.lower().endswith('.kmz'):
                import zipfile
                with zipfile.ZipFile(file_path, 'r') as kmz:
                    # Look for doc.kml or any .kml file
                    kml_files = [f for f in kmz.namelist() if f.lower().endswith('.kml')]
                    if kml_files:
                        content = kmz.read(kml_files[0]).decode('utf-8', errors='ignore')
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            
            # Parse XML and extract coordinates
            try:
                root = ET.fromstring(content)
                
                # Extract all coordinate elements
                for coord_elem in root.iter():
                    if 'coordinates' in coord_elem.tag.lower():
                        if coord_elem.text:
                            coordinates.extend(coord_elem.text.strip().split())
                
                # Extract placemark names and descriptions
                placemarks = []
                for placemark in root.iter():
                    if 'placemark' in placemark.tag.lower():
                        name = placemark.find('.//{http://www.opengis.net/kml/2.2}name')
                        desc = placemark.find('.//{http://www.opengis.net/kml/2.2}description')
                        if name is not None:
                            placemarks.append(name.text or '')
                        if desc is not None:
                            placemarks.append(desc.text or '')
                
                text_content = ' '.join(placemarks)
                keywords = self.extract_keywords(text_content)
                
            except ET.ParseError:
                text_content = content
                keywords = self.extract_keywords(content)
            
            # Convert coordinates to MGRS if possible
            mgrs_coords = []
            try:
                import mgrs
                converter = mgrs.MGRS()
                
                for coord_str in coordinates:
                    try:
                        # Parse lat,lon from coordinate string
                        parts = coord_str.split(',')
                        if len(parts) >= 2:
                            lon, lat = float(parts[0]), float(parts[1])
                            mgrs_coord = converter.toMGRS(lat, lon)
                            mgrs_coords.append(mgrs_coord)
                    except:
                        continue
            except ImportError:
                pass
            
            return {
                'full_text': text_content[:10000],
                'mgrs_coordinates': mgrs_coords,
                'gps_coordinates': coordinates[:20],  # Limit to first 20
                'keywords': ' '.join(keywords),
                'file_size': os.path.getsize(file_path)
            }
        except Exception as e:
            raise Exception(f"Error processing KML file: {str(e)}")
    
    def process_file(self, file_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """Process a single file and extract metadata"""
        if options is None:
            options = {'extractText': True, 'extractCoordinates': True}
        
        file_path = Path(file_path)
        file_stat = file_path.stat()
        file_hash = self.calculate_file_hash(str(file_path))
        
        # Basic metadata
        result = {
            'id': file_hash,
            'file_hash': file_hash,
            'highest_classification': 'UNCLASSIFIED',  # Default
            'caveats': '',
            'file_path': str(file_path),
            'locations': '',
            'timeframes': '',
            'subjects': file_path.stem,  # Use filename as subject
            'topics': self.determine_topic(file_path),
            'keywords': '',
            'MGRS': '',
            'images': None,
            'full_text': '',
            'processed_time': datetime.now().isoformat()
        }
        
        try:
            # Process based on file type
            ext = file_path.suffix.lower()
            
            if ext in self.SUPPORTED_EXTENSIONS['text']:
                data = self.process_text_file(str(file_path))
            elif ext in self.SUPPORTED_EXTENSIONS['kml']:
                data = self.process_kml_file(str(file_path))
            elif ext in self.SUPPORTED_EXTENSIONS['pdf']:
                data = self.process_pdf_file(str(file_path))
            elif ext == '.docx':
                data = self.process_word_file(str(file_path))
            elif ext == '.doc':
                data = self.process_legacy_word_file(str(file_path))
            elif ext in self.SUPPORTED_EXTENSIONS['excel']:
                data = self.process_excel_file(str(file_path))
            elif ext in self.SUPPORTED_EXTENSIONS['powerpoint']:
                data = self.process_powerpoint_file(str(file_path))
            else:
                # Generic processing
                data = {'full_text': f"File: {file_path.name}", 'keywords': file_path.stem}
            
            # Update result with extracted data
            if options.get('extractText', True):
                result['full_text'] = data.get('full_text', '')
                result['keywords'] = data.get('keywords', '')
            
            if options.get('extractCoordinates', True):
                mgrs_coords = data.get('mgrs_coordinates', [])
                if mgrs_coords:
                    result['MGRS'] = '|'.join(mgrs_coords)
                    self.stats['coordinates_found'] += len(mgrs_coords)
            
            self.stats['files_processed'] += 1
            self.stats['total_size'] += data.get('file_size', 0)
            
        except Exception as e:
            print(f"Error processing {file_path}: {str(e)}")
            result['full_text'] = f"Error processing file: {str(e)}"
            self.stats['files_failed'] += 1
        
        return result
    
    def process_pdf_file(self, file_path: str) -> Dict[str, Any]:
        """Process PDF files"""
        try:
            import PyPDF2
            text_content = ""
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text_content += page.extract_text() + "\n"
            
            mgrs_coords, gps_coords = self.extract_coordinates(text_content)
            keywords = self.extract_keywords(text_content)
            
            return {
                'full_text': text_content[:10000],  # Limit to first 10k chars
                'mgrs_coordinates': mgrs_coords,
                'gps_coordinates': gps_coords,
                'keywords': ' '.join(keywords),
                'file_size': os.path.getsize(file_path)
            }
        except ImportError:
            return {
                'full_text': f"PDF file: {os.path.basename(file_path)} (PyPDF2 not installed)",
                'keywords': os.path.splitext(os.path.basename(file_path))[0],
                'file_size': os.path.getsize(file_path)
            }
        except Exception as e:
            return {
                'full_text': f"Error processing PDF: {str(e)}",
                'keywords': os.path.splitext(os.path.basename(file_path))[0],
                'file_size': os.path.getsize(file_path)
            }

    def process_word_file(self, file_path: str) -> Dict[str, Any]:
        """Process Word documents (.docx)"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            text_content = ""
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text_content += cell.text + " "
                    text_content += "\n"
            
            mgrs_coords, gps_coords = self.extract_coordinates(text_content)
            keywords = self.extract_keywords(text_content)
            
            return {
                'full_text': text_content[:10000],  # Limit to first 10k chars
                'mgrs_coordinates': mgrs_coords,
                'gps_coordinates': gps_coords,
                'keywords': ' '.join(keywords),
                'file_size': os.path.getsize(file_path)
            }
        except ImportError:
            return {
                'full_text': f"Word file: {os.path.basename(file_path)} (python-docx not installed)",
                'keywords': os.path.splitext(os.path.basename(file_path))[0],
                'file_size': os.path.getsize(file_path)
            }
        except Exception as e:
            return {
                'full_text': f"Error processing Word document: {str(e)}",
                'keywords': os.path.splitext(os.path.basename(file_path))[0],
                'file_size': os.path.getsize(file_path)
            }

    def process_legacy_word_file(self, file_path: str) -> Dict[str, Any]:
        """Process legacy Word documents (.doc) - basic implementation"""
        try:
            # For .doc files, we'd need python-docx2txt or similar
            # For now, just return basic info
            return {
                'full_text': f"Legacy Word file: {os.path.basename(file_path)} (Limited .doc support)",
                'keywords': os.path.splitext(os.path.basename(file_path))[0],
                'file_size': os.path.getsize(file_path)
            }
        except Exception as e:
            return {
                'full_text': f"Error processing .doc file: {str(e)}",
                'keywords': os.path.splitext(os.path.basename(file_path))[0],
                'file_size': os.path.getsize(file_path)
            }
    
    def determine_topic(self, file_path: Path) -> str:
        """Determine topic based on filename and path"""
        path_parts = str(file_path).lower()
        
        if any(word in path_parts for word in ['intel', 'intelligence', 'report']):
            return 'Intelligence'
        elif any(word in path_parts for word in ['map', 'geo', 'location', 'coord']):
            return 'Geographic'
        elif any(word in path_parts for word in ['infra', 'infrastructure', 'bridge', 'road']):
            return 'Infrastructure'
        elif any(word in path_parts for word in ['security', 'threat', 'risk']):
            return 'Security'
        else:
            return 'General'
    
    def calculate_file_hash(self, file_path: str) -> str:
        """Calculate SHA256 hash of file"""
        hash_sha256 = hashlib.sha256()
        try:
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_sha256.update(chunk)
            return hash_sha256.hexdigest()
        except Exception:
            # Fallback to path-based hash if file can't be read
            return hashlib.sha256(file_path.encode()).hexdigest()
    
    def scan_folder(self, folder_path: str, options: Dict[str, Any] = None) -> List[str]:
        """Scan folder and return list of supported files"""
        if options is None:
            options = {'recursive': True, 'fileTypes': self.get_supported_formats()}
        
        files = []
        folder_path = Path(folder_path)
        
        pattern = "**/*" if options.get('recursive', True) else "*"
        
        for file_path in folder_path.glob(pattern):
            if file_path.is_file() and self.is_supported_file(str(file_path), options.get('fileTypes')):
                files.append(str(file_path))
        
        return sorted(files)
    
    def process_excel_file(self, file_path: str) -> Dict[str, Any]:
        """Process Excel files (.xlsx, .xls)"""
        try:
            import pandas as pd
            
            # Read Excel file
            excel_file = pd.ExcelFile(file_path)
            text_content = ""
            
            # Extract text from all sheets
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text_content += f"Sheet: {sheet_name}\n"
                text_content += df.to_string() + "\n\n"
            
            mgrs_coords, gps_coords = self.extract_coordinates(text_content)
            keywords = self.extract_keywords(text_content)
            
            return {
                'full_text': text_content[:10000],  # Limit to first 10k chars
                'mgrs_coordinates': mgrs_coords,
                'gps_coordinates': gps_coords,
                'keywords': ' '.join(keywords),
                'file_size': os.path.getsize(file_path)
            }
        except ImportError:
            return {
                'full_text': f"Excel file: {os.path.basename(file_path)} (pandas not installed)",
                'keywords': os.path.splitext(os.path.basename(file_path))[0],
                'file_size': os.path.getsize(file_path)
            }
        except Exception as e:
            return {
                'full_text': f"Error processing Excel file: {str(e)}",
                'keywords': os.path.splitext(os.path.basename(file_path))[0],
                'file_size': os.path.getsize(file_path)
            }

    def process_powerpoint_file(self, file_path: str) -> Dict[str, Any]:
        """Process PowerPoint files (.pptx, .ppt)"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_content = ""
            
            # Extract text from all slides
            for i, slide in enumerate(prs.slides):
                text_content += f"Slide {i+1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
                text_content += "\n"
            
            mgrs_coords, gps_coords = self.extract_coordinates(text_content)
            keywords = self.extract_keywords(text_content)
            
            return {
                'full_text': text_content[:10000],  # Limit to first 10k chars
                'mgrs_coordinates': mgrs_coords,
                'gps_coordinates': gps_coords,
                'keywords': ' '.join(keywords),
                'file_size': os.path.getsize(file_path)
            }
        except ImportError:
            return {
                'full_text': f"PowerPoint file: {os.path.basename(file_path)} (python-pptx not installed)",
                'keywords': os.path.splitext(os.path.basename(file_path))[0],
                'file_size': os.path.getsize(file_path)
            }
        except Exception as e:
            return {
                'full_text': f"Error processing PowerPoint file: {str(e)}",
                'keywords': os.path.splitext(os.path.basename(file_path))[0],
                'file_size': os.path.getsize(file_path)
            }
